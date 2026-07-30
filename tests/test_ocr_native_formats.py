from __future__ import annotations

import io
import zipfile
from types import SimpleNamespace

import pytest

from server.ocr import formats, native
from server.ocr.classify import classify
from server.ocr import pipeline
from server.routes.upload_helpers import validate_document_upload


def test_manifest_loader_rejects_missing_group(tmp_path):
    manifest = tmp_path / "formats.json"
    manifest.write_text('{"version": 1, "groups": {"text": [".txt"]}}', encoding="utf-8")

    with pytest.raises(RuntimeError, match="missing groups"):
        formats.load_format_manifest(manifest)


def test_every_manifest_suffix_has_a_non_manual_backend_route(tmp_path):
    for suffix in formats.ALL_SUPPORTED_SUFFIXES:
        path = tmp_path / f"document{suffix}"
        if suffix == ".pdf":
            path.write_bytes(b"%PDF-1.7\n/Type /Page\n/DCTDecode")
        else:
            path.write_bytes(b"placeholder")

        result = classify(path)

        assert result["route"] in {"native", "convert", "ocr"}, suffix
        assert result["handler"] != "unknown", suffix


def _zip_package(entries: dict[str, bytes]) -> bytes:
    output = io.BytesIO()
    with zipfile.ZipFile(output, "w") as archive:
        for name, content in entries.items():
            archive.writestr(name, content)
    return output.getvalue()


def _valid_content(suffix: str) -> bytes:
    text = "招标文件有效底稿".encode()
    if suffix in formats.suffixes("text"):
        return text
    if suffix == ".pdf":
        return b"%PDF-1.7\n/Type /Page\n/DCTDecode"
    if suffix in {".doc", ".xls", ".ppt"}:
        return b"\xd0\xcf\x11\xe0\xa1\xb1\x1a\xe1" + b"legacy-office"
    if suffix == ".docx":
        return _zip_package(
            {"[Content_Types].xml": b"types", "word/document.xml": b"<w:t>Tender text</w:t>"}
        )
    if suffix in {".xlsx", ".xlsm"}:
        return _zip_package(
            {"[Content_Types].xml": b"types", "xl/workbook.xml": b"<workbook/>"}
        )
    if suffix == ".xlsb":
        return _zip_package(
            {"[Content_Types].xml": b"types", "xl/workbook.bin": b"workbook"}
        )
    if suffix == ".pptx":
        return _zip_package(
            {"[Content_Types].xml": b"types", "ppt/presentation.xml": b"<presentation/>"}
        )
    if suffix in {".odt", ".ods", ".odp"}:
        subtype = {".odt": "text", ".ods": "spreadsheet", ".odp": "presentation"}[suffix]
        return _zip_package(
            {"mimetype": f"application/vnd.oasis.opendocument.{subtype}".encode()}
        )
    return {
        ".png": b"\x89PNG\r\n\x1a\nimage",
        ".jpg": b"\xff\xd8\xffimage",
        ".jpeg": b"\xff\xd8\xffimage",
        ".tif": b"II*\x00image",
        ".tiff": b"MM\x00*image",
        ".bmp": b"BMimage",
        ".webp": b"RIFF\x04\x00\x00\x00WEBPimage",
    }[suffix]


def test_every_manifest_suffix_uploads_routes_and_yields_nonempty_contract(tmp_path, monkeypatch):
    monkeypatch.setattr(
        pipeline,
        "_call_native_read",
        lambda *_args, **_kwargs: {"kind": "text", "blocks": ["native 底稿"]},
    )
    monkeypatch.setattr(
        pipeline,
        "_convert_and_dispatch",
        lambda *_args, **_kwargs: {"kind": "pdf_text", "route": "convert", "blocks": ["convert 底稿"]},
    )
    monkeypatch.setattr(
        pipeline,
        "_call_recognize_with_seal",
        lambda *_args, **_kwargs: {
            "kind": "ocr",
            "route": "ocr",
            "pages": [{"page_number": 1, "markdown": "OCR 底稿"}],
        },
    )

    for suffix in sorted(formats.ALL_SUPPORTED_SUFFIXES):
        content = _valid_content(suffix)
        validate_document_upload(f"document{suffix}", content)
        path = tmp_path / f"document{suffix}"
        path.write_bytes(content)

        route = classify(path)
        result = pipeline._dispatch_extract(
            path, run_seal=False, purpose="contract-test", on_page=None
        )
        rendered = pipeline._render_body(result)

        assert route["route"] in {"native", "convert", "ocr"}, suffix
        assert result.get("kind") != "manual", suffix
        assert rendered.strip(), suffix
        assert pipeline.is_ocr_text_valid(f"### 文件: document{suffix}\n{rendered}"), suffix


@pytest.mark.parametrize(
    ("suffix", "route", "handler"),
    [
        (".xls", "native", "excel_xls"),
        (".xlsb", "native", "excel_xlsb"),
        (".pptx", "native", "presentation"),
        (".ppt", "convert", "office_convert"),
        (".odt", "convert", "office_convert"),
        (".ods", "convert", "office_convert"),
        (".odp", "convert", "office_convert"),
    ],
)
def test_classify_extended_office_matrix(tmp_path, suffix, route, handler):
    path = tmp_path / f"document{suffix}"
    path.write_bytes(b"placeholder")

    result = classify(path)

    assert result["route"] == route
    assert result["handler"] == handler


def test_short_docx_with_real_text_stays_native(tmp_path):
    import zipfile

    path = tmp_path / "short.docx"
    with zipfile.ZipFile(path, "w") as archive:
        archive.writestr("word/document.xml", "<w:p><w:t>采购公告</w:t></w:p>")

    result = classify(path)

    assert result["route"] == "native"
    assert result["handler"] == "word"


def test_scanned_docx_routes_to_office_conversion(tmp_path):
    import zipfile

    path = tmp_path / "scan.docx"
    with zipfile.ZipFile(path, "w") as archive:
        archive.writestr("word/document.xml", "<w:p></w:p>")
        archive.writestr("word/media/image1.png", b"\x89PNG")

    result = classify(path)

    assert result["route"] == "convert"
    assert result["handler"] == "office_convert"


def test_xls_reader_never_imports_openpyxl(monkeypatch, tmp_path):
    sheet = SimpleNamespace(
        name="报价",
        nrows=2,
        row_values=lambda index: [["项目", "金额"], ["服务", 12.5]][index],
    )
    workbook = SimpleNamespace(sheets=lambda: [sheet], release_resources=lambda: None)
    imports: list[str] = []

    def fake_require(module: str, package: str):
        imports.append(module)
        assert module == "xlrd"
        return SimpleNamespace(open_workbook=lambda *args, **kwargs: workbook)

    monkeypatch.setattr(native, "_require", fake_require)
    path = tmp_path / "报价.xls"
    path.write_bytes(b"xls")

    result = native.read_excel_xls(path)

    assert imports == ["xlrd"]
    assert result["tables"][0]["rows"][1] == ["服务", "12.5"]


def test_xlsb_reader_extracts_rows(monkeypatch, tmp_path):
    rows = [
        [SimpleNamespace(v="项目"), SimpleNamespace(v="金额")],
        [SimpleNamespace(v="服务"), SimpleNamespace(v=8)],
    ]

    class Sheet:
        def __enter__(self):
            return iter(rows)

        def __exit__(self, *_args):
            return False

    class Workbook:
        sheets = ["报价"]

        def get_sheet(self, _name):
            return Sheet()

        def __enter__(self):
            return self

        def __exit__(self, *_args):
            return False

    monkeypatch.setattr(
        native,
        "_require",
        lambda module, package: SimpleNamespace(open_workbook=lambda _path: Workbook()),
    )

    result = native.read_excel_xlsb(tmp_path / "报价.xlsb")

    assert result["tables"][0]["rows"][1] == ["服务", "8"]


def test_pptx_reader_extracts_shape_and_table_text(monkeypatch, tmp_path):
    text_shape = SimpleNamespace(
        has_text_frame=True, text="投标方案", has_table=False, shape_type=1
    )
    table = SimpleNamespace(
        rows=[SimpleNamespace(cells=[SimpleNamespace(text="工期"), SimpleNamespace(text="30天")])]
    )
    table_shape = SimpleNamespace(
        has_text_frame=False, has_table=True, table=table, shape_type=19
    )
    image_shape = SimpleNamespace(
        has_text_frame=False, has_table=False, shape_type=13
    )
    presentation = SimpleNamespace(
        slides=[SimpleNamespace(shapes=[text_shape, table_shape, image_shape])]
    )
    monkeypatch.setattr(
        native,
        "_require",
        lambda module, package: SimpleNamespace(Presentation=lambda _path: presentation),
    )

    result = native.read_presentation(tmp_path / "方案.pptx")

    assert result["blocks"] == ["投标方案"]
    assert result["tables"][0]["rows"] == [["工期", "30天"]]
    assert result["image_count"] == 1
    assert result["text_char_count"] == len("投标方案工期30天")


def test_real_pptx_reader_detects_scanned_picture_with_short_title(tmp_path):
    from PIL import Image
    from pptx import Presentation
    from pptx.util import Inches

    scan = tmp_path / "scan.png"
    Image.new("RGB", (120, 80), "white").save(scan)
    path = tmp_path / "扫描型投标方案.pptx"
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    title = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(4), Inches(0.5))
    title.text = "投标方案"
    slide.shapes.add_picture(str(scan), Inches(0.5), Inches(1), width=Inches(4))
    presentation.save(path)

    result = native.read_presentation(path)

    assert result["blocks"] == ["投标方案"]
    assert result["image_count"] == 1
    assert result["text_char_count"] == len("投标方案")


def test_real_grouped_pptx_picture_and_short_title_trigger_scan_upgrade(tmp_path):
    from PIL import Image
    from pptx import Presentation
    from pptx.util import Inches

    from server.ocr.pipeline import _presentation_needs_ocr

    scan = tmp_path / "grouped-scan.png"
    Image.new("RGB", (120, 80), "white").save(scan)
    path = tmp_path / "组合扫描型方案.pptx"
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    group = slide.shapes.add_group_shape()
    title = group.shapes.add_textbox(Inches(0.2), Inches(0.2), Inches(3), Inches(0.5))
    title.text = "投标方案"
    group.shapes.add_picture(str(scan), Inches(0.2), Inches(1), width=Inches(3))
    presentation.save(path)

    result = native.read_presentation(path)

    assert result["blocks"] == ["投标方案"]
    assert result["image_count"] == 1
    assert result["text_char_count"] == len("投标方案")
    assert _presentation_needs_ocr(result) is True


def test_nested_group_text_table_and_picture_are_counted_once(monkeypatch, tmp_path):
    from server.ocr.pipeline import _presentation_needs_ocr

    long_text = "x" * 80
    text_shape = SimpleNamespace(
        has_text_frame=True, text=long_text, has_table=False, shape_type=1
    )
    table = SimpleNamespace(
        rows=[SimpleNamespace(cells=[SimpleNamespace(text="工期"), SimpleNamespace(text="30天")])]
    )
    table_shape = SimpleNamespace(
        has_text_frame=False, has_table=True, table=table, shape_type=19
    )
    image_shape = SimpleNamespace(has_text_frame=False, has_table=False, shape_type=13)
    inner_group = SimpleNamespace(
        has_text_frame=False,
        has_table=False,
        shape_type=6,
        shapes=[text_shape, table_shape, image_shape],
    )
    outer_group = SimpleNamespace(
        has_text_frame=False,
        has_table=False,
        shape_type=6,
        shapes=[inner_group],
    )
    presentation = SimpleNamespace(slides=[SimpleNamespace(shapes=[outer_group])])
    monkeypatch.setattr(
        native,
        "_require",
        lambda module, package: SimpleNamespace(Presentation=lambda _path: presentation),
    )

    result = native.read_presentation(tmp_path / "nested.pptx")

    assert result["blocks"] == [long_text]
    assert result["tables"] == [{"rows": [["工期", "30天"]]}]
    assert result["image_count"] == 1
    assert result["text_char_count"] == len(long_text + "工期30天")
    assert _presentation_needs_ocr(result) is False
